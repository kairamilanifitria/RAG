import logging
import json
import os
from pathlib import Path
import warnings
import torch
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import hashlib
import re
from bs4 import BeautifulSoup
import html2text


os.environ["HF_HUB_DISABLE_HARDLINKS"] = "1"
os.environ["HF_HUB_DISABLE_SYMLINKS"] = "1" 

# Ensure PyTorch uses GPU if available
torch.set_default_device("cuda" if torch.cuda.is_available() else "cpu")

print("CUDA Available:", torch.cuda.is_available())
print("Using Device:", torch.cuda.get_device_name(0) if torch.cuda.is_available() else "CPU")


from docling.datamodel.base_models import InputFormat
from docling_core.types.doc import ImageRefMode
from docling.document_converter import DocumentConverter, PdfFormatOption, WordFormatOption
from docling.pipeline.standard_pdf_pipeline import StandardPdfPipeline
from docling.datamodel.pipeline_options import PdfPipelineOptions
from docling.pipeline.simple_pipeline import SimplePipeline
from docling.datamodel.settings import settings

warnings.filterwarnings("ignore")

# Logging setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

IMAGE_RESOLUTION_SCALE = 2.0

INPUT_DIR = r"E:\\CODE\\RAG\\Documents"  # Folder containing PDFs/DOCs
OUTPUT_DIR = r"E:\\CODE\\RAG\\Documents\\output"   # Folder for Markdown outputs

os.makedirs(INPUT_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

def create_pipeline_options(input_format):
    if input_format == InputFormat.PDF:
        return PdfFormatOption(
            pipeline_options=PdfPipelineOptions(
                do_table_structure=True,
                generate_page_images=True,
                generate_picture_images=True,
                images_scale=IMAGE_RESOLUTION_SCALE,
            )
        )
    elif input_format == InputFormat.DOCX:
        return WordFormatOption(pipeline_cls=SimplePipeline)
    return None


def initialize_converter():
    allowed_formats = [InputFormat.PDF, InputFormat.DOCX]
    format_options = {
        fmt: create_pipeline_options(fmt)
        for fmt in allowed_formats if create_pipeline_options(fmt)
    }
    return DocumentConverter(allowed_formats=allowed_formats, format_options=format_options)

# Modify this to allow PDF, DOCX, CSV, XLSX, PPT
def convert_and_save():
    """Converts documents to Markdown and saves the output."""
    input_paths = list(Path(INPUT_DIR).glob("*.*"))

    if not input_paths:
        logger.warning("No input files found in the directory.")
        return
    
    # Docling handles PDF + DOCX
    converter = initialize_converter()
    docling_paths = [p for p in input_paths if p.suffix.lower() in [".pdf", ".docx"]]
    other_paths = [p for p in input_paths if p.suffix.lower() in [".csv", ".xlsx", ".txt", ".pptx", ".json", ".html", ".md"]]

    # Run Docling
    if docling_paths:
        results = converter.convert_all(docling_paths)
        for res in results:
            file_name = res.input.file.stem
            md_path = Path(OUTPUT_DIR) / f"{file_name}.md"
            res.document.save_as_markdown(md_path, image_mode=ImageRefMode.REFERENCED)
            logger.info(f"Markdown saved to {md_path}")

    # Run fallback converters
    for file_path in other_paths:
        file_name = file_path.stem
        md_path = Path(OUTPUT_DIR) / f"{file_name}.md"
        convert_non_docling(file_path, md_path)
        logger.info(f"Markdown saved to {md_path} (fallback parser)")


def convert_non_docling(file_path, output_path):
    ext = file_path.suffix.lower()

    if ext == ".csv":
        df = pd.read_csv(file_path)
        md_parts = []
        for i, row in df.iterrows():
            md_parts.append(f"## Data {i + 1}")
            for col, val in row.items():
                md_parts.append(f"- {col}: {val}")
            md_parts.append("")  # spacing
        md = "\n".join(md_parts)

    elif ext == ".xlsx":
        xls = pd.ExcelFile(file_path)
        md_parts = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            # md_parts.append(f"# Sheet name: {sheet_name}\n")
            for i, row in df.iterrows():
                md_parts.append(f"## Data {i + 1}")
                for col, val in row.items():
                    md_parts.append(f"- {col}: {val}")
                md_parts.append("")  # spacing
        md = "\n".join(md_parts)

    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read().strip()

        paragraphs = [p.strip() for p in re.split(r'\n\s*\n', content) if p.strip()]

        md_lines = []
        for i, para in enumerate(paragraphs, start=1):
            md_lines.append(f"## Paragraph {i}")
            md_lines.append(para)
            md_lines.append("")  # blank line between sections

        md = "\n".join(md_lines)

    elif ext == ".pptx":
        md = parse_ppt(file_path)

    elif ext == ".json":
        md = parse_json(file_path)

    elif ext == ".md":
        with open(file_path, "r", encoding="utf-8") as f:
            md = f.read()

    elif ext == ".html":
        md = parse_html(file_path)

    else:
        return  # unsupported

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(md)

def extract_images_from_shape(shape, artifact_dir, image_counter):
    """Recursively extract images from shapes (including groups)."""
    images = []
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        image_bytes = image.blob
        image_ext = image.ext or "png"

        # Use hash for unique filenames
        hash_digest = hashlib.sha256(image_bytes).hexdigest()[:12]
        image_filename = f"image_{image_counter:06d}_{hash_digest}.{image_ext}"
        image_path = artifact_dir / image_filename

        with open(image_path, "wb") as f:
            f.write(image_bytes)

        images.append((image_path, image_counter + 1))

    # If shape is a group, recursively check its sub-shapes
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            sub_images = extract_images_from_shape(sub_shape, artifact_dir, image_counter + len(images))
            images.extend(sub_images)

    return images

def parse_ppt(file_path):
    prs = Presentation(file_path)
    artifact_dir = Path(OUTPUT_DIR) / f"{Path(file_path).stem}_artifacts"
    artifact_dir.mkdir(parents=True, exist_ok=True)

    slides = []
    image_counter = 0

    for i, slide in enumerate(prs.slides, start=1):
        slide_content = [f"## Slide {i}"]

        for shape in slide.shapes:
            # Extract text
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())

            # Extract images (including nested/grouped)
            extracted = extract_images_from_shape(shape, artifact_dir, image_counter)
            for img_path, new_counter in extracted:
                slide_content.append(f"![Image]({img_path})")
                image_counter = new_counter

        slides.append("\n".join(slide_content))

    md = "\n\n".join(slides)
    return md

def parse_json(file_path):
    """
    Parse a JSON file into Markdown format.
    Converts keys and values into readable sections.
    Handles both dicts and lists.
    """
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def format_json(data, level=2):
        """Recursively format JSON as Markdown text."""
        md_lines = []
        if isinstance(data, dict):
            for key, value in data.items():
                header = "#" * level + f" {key}"
                md_lines.append(header)
                md_lines.append("")
                md_lines.append(format_json(value, level + 1))
        elif isinstance(data, list):
            for i, item in enumerate(data, start=1):
                md_lines.append(f"- **Item {i}:**")
                md_lines.append(format_json(item, level + 1))
        else:
            md_lines.append(str(data))
        return "\n".join(md_lines)

    md = format_json(data)
    return md

def parse_html(file_path):
    """
    Parse an HTML file and convert it to Markdown.
    Keeps text, headings, links, and images in readable Markdown format.
    """
    with open(file_path, "r", encoding="utf-8") as f:
        html_content = f.read()

    h = html2text.HTML2Text()
    h.ignore_links = False
    h.ignore_images = False
    h.ignore_emphasis = False
    h.bypass_tables = False
    h.body_width = 0  # Prevent wrapping

    md = h.handle(html_content).strip()

    # Optional: cleanup empty lines
    md = "\n".join([line.strip() for line in md.splitlines() if line.strip()])

    return md

def main():
    settings.debug.profile_pipeline_timings = True
    convert_and_save()

if __name__ == "__main__":
    main()